#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT Processor - EZTeach Skill
PPT内容提取、动画窗格解析、自动转PDF全流程处理
"""

import os
import json
from typing import Dict, List, Any, Optional
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import win32com.client
    WIN32_AVAILABLE = True
except ImportError:
    WIN32_AVAILABLE = False


class PPTProcessor:
    """PPT处理器类"""
    
    def __init__(self, ppt_path: str):
        """
        初始化PPT处理器
        
        Args:
            ppt_path: PPT文件路径
        """
        self.ppt_path = Path(ppt_path)
        self.prs = None
        self.slides_data = []
        self.animations_data = []
        
        if PPTX_AVAILABLE and self.ppt_path.exists():
            self.prs = Presentation(str(self.ppt_path))
    
    def extract_slide_content(self, slide) -> Dict[str, Any]:
        """
        提取单页幻灯片内容
        
        Args:
            slide: pptx Slide对象
            
        Returns:
            幻灯片内容字典
        """
        content = {
            "slide_index": slide.slide_id,
            "title": "",
            "texts": [],
            "shapes": [],
            "images": [],
            "tables": [],
            "notes": ""
        }
        
        # 提取标题
        if slide.shapes.title:
            content["title"] = slide.shapes.title.text
        
        # 提取所有形状和文本
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": shape.shape_type,
                "has_text": shape.has_text_frame,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
            
            # 提取文本
            if shape.has_text_frame:
                text_content = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        text_content.append({
                            "text": para_text,
                            "level": paragraph.level,
                            "font_size": paragraph.font.size.pt if paragraph.font.size else None
                        })
                        content["texts"].append(para_text)
                shape_info["text_content"] = text_content
            
            # 提取图片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_info = {
                    "name": shape.name,
                    "ext": shape.image.ext,
                    "size": shape.image.size
                }
                content["images"].append(image_info)
                shape_info["is_image"] = True
            
            # 提取表格
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                content["tables"].append(table_data)
                shape_info["is_table"] = True
                shape_info["table_data"] = table_data
            
            content["shapes"].append(shape_info)
        
        # 提取备注
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                content["notes"] = notes_slide.notes_text_frame.text
        
        return content
    
    def extract_all_content(self) -> Dict[str, Any]:
        """
        提取所有幻灯片内容
        
        Returns:
            完整PPT内容字典
        """
        if not self.prs:
            return {"error": "PPT文件无法加载或python-pptx不可用"}
        
        result = {
            "file_name": self.ppt_path.name,
            "total_slides": len(self.prs.slides),
            "slide_width": self.prs.slide_width,
            "slide_height": self.prs.slide_height,
            "slides": []
        }
        
        for i, slide in enumerate(self.prs.slides):
            slide_content = self.extract_slide_content(slide)
            slide_content["slide_number"] = i + 1
            result["slides"].append(slide_content)
        
        self.slides_data = result
        return result
    
    def parse_animations(self) -> List[Dict[str, Any]]:
        """
        解析动画窗格信息
        
        Returns:
            动画信息列表
        """
        animations = []
        
        if not self.prs:
            return animations
        
        # python-pptx对动画支持有限，这里提取基本的时序信息
        for i, slide in enumerate(self.prs.slides):
            slide_animations = {
                "slide_number": i + 1,
                "has_timeline": hasattr(slide, 'timeline'),
                "animations_count": 0,
                "effects": []
            }
            
            # 尝试提取动画效果
            try:
                if hasattr(slide, 'timeline') and slide.timeline:
                    timeline = slide.timeline
                    if hasattr(timeline, 'mainSequence'):
                        seq = timeline.mainSequence
                        slide_animations["animations_count"] = len(seq)
                        for effect in seq:
                            effect_info = {
                                "type": str(type(effect)),
                                "display_name": getattr(effect, 'displayName', ''),
                                "trigger": getattr(effect, 'trigger', '')
                            }
                            slide_animations["effects"].append(effect_info)
            except Exception as e:
                slide_animations["parse_error"] = str(e)
            
            animations.append(slide_animations)
        
        self.animations_data = animations
        return animations
    
    def export_to_json(self, output_path: str) -> str:
        """
        导出PPT内容为JSON
        
        Args:
            output_path: 输出文件路径
            
        Returns:
            输出文件路径
        """
        content = self.extract_all_content()
        animations = self.parse_animations()
        content["animations"] = animations
        
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(content, f, ensure_ascii=False, indent=2)
        
        return str(output_path)
    
    def convert_to_pdf(self, output_path: Optional[str] = None) -> str:
        """
        PPT转换为PDF（Windows环境使用COM，其他环境提示）
        
        Args:
            output_path: 输出PDF路径
            
        Returns:
            PDF文件路径或提示信息
        """
        if output_path is None:
            output_path = self.ppt_path.with_suffix('.pdf')
        else:
            output_path = Path(output_path)
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Windows环境使用win32com
        if WIN32_AVAILABLE:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                powerpoint.Visible = 1
                
                deck = powerpoint.Presentations.Open(str(self.ppt_path.absolute()))
                deck.SaveAs(str(output_path.absolute()), 32)  # 32 = PDF
                deck.Close()
                powerpoint.Quit()
                
                return str(output_path)
            except Exception as e:
                return f"PDF转换失败: {str(e)}"
        else:
            return f"PDF转换需要Windows环境和pywin32库。请手动导出或使用LibreOffice: libreoffice --headless --convert-to pdf {self.ppt_path}"
    
    def get_teaching_outline(self) -> List[str]:
        """
        提取教学大纲（从标题和一级文本）
        
        Returns:
            大纲列表
        """
        outline = []
        
        if not self.slides_data:
            self.extract_all_content()
        
        for slide in self.slides_data.get("slides", []):
            if slide.get("title"):
                outline.append(slide["title"])
        
        return outline


def process_ppt_file(ppt_path: str, output_dir: Optional[str] = None) -> Dict[str, Any]:
    """
    处理PPT文件的主函数
    
    Args:
        ppt_path: PPT文件路径
        output_dir: 输出目录
        
    Returns:
        处理结果字典
    """
    processor = PPTProcessor(ppt_path)
    
    result = {
        "success": True,
        "ppt_path": ppt_path,
        "content_extracted": False,
        "json_exported": None,
        "pdf_converted": None,
        "outline": []
    }
    
    # 提取内容
    content = processor.extract_all_content()
    if "error" not in content:
        result["content_extracted"] = True
        result["total_slides"] = content.get("total_slides", 0)
        result["outline"] = processor.get_teaching_outline()
    
    # 导出JSON
    if output_dir:
        output_dir = Path(output_dir)
        json_path = output_dir / f"{Path(ppt_path).stem}_content.json"
        result["json_exported"] = processor.export_to_json(str(json_path))
    
    # 转换PDF
    if output_dir:
        pdf_path = output_dir / f"{Path(ppt_path).stem}.pdf"
        result["pdf_converted"] = processor.convert_to_pdf(str(pdf_path))
    
    return result


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        ppt_file = sys.argv[1]
        output = sys.argv[2] if len(sys.argv) > 2 else None
        result = process_ppt_file(ppt_file, output)
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        print("Usage: python ppt_processor.py <ppt_file> [output_dir]")
